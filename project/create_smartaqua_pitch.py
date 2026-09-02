from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.dml import MSO_THEME_COLOR

OUT = r"C:\Users\HP\Desktop\project\SmartAqua_AASW2026_Pitch_Deck.pptx"
IMG = r"C:\Users\HP\Downloads\Telegram Desktop\photo_2026-07-16_11-01-16.jpg"
LOGO = r"C:\Users\HP\Desktop\project\mobile\assets\images\logo.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(8, 48, 38)
DEEP = RGBColor(5, 34, 27)
CREAM = RGBColor(247, 244, 232)
MINT = RGBColor(190, 232, 190)
GOLD = RGBColor(245, 190, 72)
MUTED = RGBColor(178, 201, 190)
WHITE = RGBColor(255, 255, 255)

def textbox(slide, text, x, y, w, h, size=18, color=CREAM, bold=False,
            font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    return shape

def logo(slide, x=0.55, y=0.35, dark=False):
    slide.shapes.add_picture(LOGO, Inches(x), Inches(y-0.10), height=Inches(0.48))
    c = BG if dark else CREAM
    textbox(slide, "SMART", x+0.58, y, 0.72, 0.3, 12, c, True)
    textbox(slide, "AQUA", x+1.26, y, 0.72, 0.3, 12, GOLD, True)

def base(title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    logo(s)
    if kicker: textbox(s, kicker.upper(), 10.1, 0.38, 2.65, 0.25, 10, GOLD, True, align=PP_ALIGN.RIGHT)
    textbox(s, title, 0.6, 1.05, 12.1, 0.65, 31, CREAM, True)
    rect(s, 0.6, 1.85, 1.0, 0.07, GOLD)
    return s

def bullets(slide, items, x, y, w, h, size=18, color=CREAM, gap=0.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    for i, item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.level=0; p.font.name="Aptos"; p.font.size=Pt(size); p.font.color.rgb=color; p.space_after=Pt(gap*10); p.bullet=True
    return box

# 1 cover
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=DEEP
rect(s, 8.45, 0, 4.88, 7.5, BG)
logo(s, 0.7, 0.55)
textbox(s, "Intelligent feeding\nfor resilient aquaculture", 0.75, 1.55, 7.4, 1.7, 34, CREAM, True)
textbox(s, "SmartAqua", 0.78, 3.55, 3.2, 0.55, 25, GOLD, True)
textbox(s, "Africa Agri-future Innovation Challenge 2026", 0.8, 4.25, 5.6, 0.35, 16, MINT, True)
textbox(s, "Ogunbunmi Philip Oluwadarasimi\nFederal University of Technology, Akure\n08169150242  •  ogunbunmiphilip@gmail.com", 0.8, 5.55, 6.4, 0.9, 14, MUTED)
s.shapes.add_picture(IMG, Inches(8.75), Inches(1.3), width=Inches(4.1), height=Inches(4.9))
textbox(s, "FIELD-TESTED PILOT", 9.1, 6.55, 3.5, 0.3, 12, GOLD, True, align=PP_ALIGN.CENTER)

# 2 problem
s=base("The aquaculture problem", "Why it matters")
bullets(s,["Inconsistent feeding increases waste and operating cost.","Farmers cannot always be present for every feeding event.","Water-temperature changes affect fish metabolism and feed demand.","Rural farms need solutions that can become independent of unreliable grid power.","Limited farm data makes it difficult to improve feeding decisions."],0.8,2.35,7.0,3.7,20)
rect(s,8.4,2.4,3.9,3.7,DEEP,True,line=MINT)
textbox(s,"The opportunity",8.8,2.8,3.0,0.4,22,GOLD,True)
textbox(s,"Make feeding consistent, measurable and responsive to real farm conditions—without requiring the farmer to be physically present.",8.8,3.55,3.0,1.45,17,CREAM)
textbox(s,"Target users: small and medium-scale African fish farms",8.8,5.45,3.0,0.52,12,MINT,True)

# 3 solution
s=base("One platform, three layers", "The solution")
for x,title,body in [(0.8,"Feeder","Motorized auger dispenser with scheduled and manual control."),(4.55,"Intelligence","Q10 temperature adjustment, feeding logs and operational sensing."),(8.3,"Farmer app","Remote control, schedules, history, alerts and device monitoring.")]:
    rect(s,x,2.25,3.25,3.25,DEEP,True,line=MINT); textbox(s,title,x+0.25,2.65,2.75,0.45,22,GOLD,True); textbox(s,body,x+0.25,3.45,2.7,1.3,18,CREAM)
textbox(s,"Current validated path",0.9,6.05,2.3,0.3,14,MINT,True)
textbox(s,"Schedule or manual request  →  Q10 adjustment  →  dispensing  →  experiment log",3.1,6.0,8.8,0.35,16,CREAM)

# 4 design
s=base("Experimental feeder design", "Prototype in the field")
s.shapes.add_picture(IMG, Inches(0.7), Inches(2.2), width=Inches(7.1), height=Inches(4.65))
rect(s,8.2,2.25,4.2,3.9,DEEP,True,line=MINT)
textbox(s,"Design characteristics",8.55,2.6,3.3,0.4,21,GOLD,True)
bullets(s,["Hopper and auger-based dispensing","Stepper motor control","Designed for repeatable quantities","Mounting frame for pond-side installation","Expandable sensor and power architecture"],8.55,3.35,3.35,2.1,16)
textbox(s,"Experimental-stage hardware shown",8.55,5.7,3.2,0.35,13,MUTED)

# 5 validated experiment
s=base("Evidence from the pilot experiment", "What we measured")
textbox(s,"56 days  •  224 feeding records  •  Manual control vs Smart Assisted",0.8,2.12,11.5,0.35,17,MINT,True)
chart_data=CategoryChartData(); chart_data.categories=["Biomass gain (g)","Final biomass (g)","ADG (g/fish/day)"]
chart_data.add_series("Manual control",(2755.33,8987,2.624)); chart_data.add_series("Smart Assisted",(3169.22,9404,3.391))
chart=s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,Inches(0.7),Inches(2.7),Inches(7.3),Inches(3.65),chart_data).chart
chart.has_legend=True; chart.legend.include_in_layout=False; chart.legend.position=2; chart.value_axis.has_major_gridlines=True
chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb=MUTED
chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb=GOLD
chart.category_axis.tick_labels.font.color.rgb = CREAM
chart.category_axis.tick_labels.font.size = Pt(12)
chart.value_axis.tick_labels.font.color.rgb = MUTED
chart.value_axis.tick_labels.font.size = Pt(11)
chart.legend.font.color.rgb = CREAM
rect(s,8.45,2.7,3.7,3.65,DEEP,True,line=MINT)
textbox(s,"Pilot signals",8.8,3.0,2.8,0.35,21,GOLD,True)
bullets(s,["+15.0% biomass gain","+29.2% average daily gain","+9.7% LER","~8.9% lower calculated FCR","Q10 factor logged across temperature changes"],8.8,3.75,2.9,2.1,17)

# 6 how q10 works
s=base("Q10 makes feeding responsive", "Core innovation")
textbox(s,"Temperature changes influence fish metabolism. SmartAqua uses a Q10 factor to adjust the requested feed quantity instead of applying one fixed amount every day.",0.8,2.25,7.0,1.1,21,CREAM)
rect(s,0.85,3.85,6.6,1.35,DEEP,True,line=MINT)
textbox(s,"Base feed quantity  ×  Q10 adjustment factor  =  temperature-responsive quantity",1.1,4.25,6.0,0.5,23,GOLD,True,align=PP_ALIGN.CENTER)
rect(s,8.15,2.3,4.1,3.8,DEEP,True,line=MINT)
textbox(s,"Observed in the trial",8.5,2.7,3.3,0.35,21,GOLD,True)
textbox(s,"Temperature range\n23.06–27.06 °C\n\nQ10 factor range\n0.866–1.165\n\nAverage dispensing deviation\n≈ 1.26 g below request",8.5,3.45,3.2,2.1,18,CREAM)

# 7 roadmap
s=base("From validated pilot to resilient platform", "Next phase")
for x,title,body in [(0.75,"Now validated","Manual feeding\nScheduled feeding\nQ10 adjustment\nTwo-month field use"),(4.45,"Next validation","Solar + battery operation\nDO and pH sensing\nComputer-vision verification\nReliability under poor connectivity"),(8.15,"Scale pathway","More farms and fish populations\nMeasured labor and feed savings\nDeployment partnerships\nCommercial productization")]:
    rect(s,x,2.35,3.45,3.55,DEEP,True,line=MINT); textbox(s,title,x+0.3,2.75,2.8,0.4,21,GOLD,True); textbox(s,body,x+0.3,3.55,2.8,2.05,16,CREAM)
textbox(s,"Advanced AI, water-quality sensing and solar operation are presented as the next validation phase—not as completed outcomes.",0.85,6.25,11.8,0.35,14,MINT,True,align=PP_ALIGN.CENTER)

# 8 ask
s=base("What support will unlock", "The ask")
bullets(s,["Mentorship in aquaculture trials, product validation and commercialization.","Pilot partnerships with additional fish farms and research institutions.","Technical support for solar power, water-quality sensing and computer vision.","Access to stakeholders, investors and African aquaculture networks.","Funding to expand the field trial and quantify feed efficiency, labor savings and fish growth."],0.8,2.35,7.6,3.5,19)
rect(s,8.85,2.5,3.3,3.55,DEEP,True,line=GOLD)
textbox(s,"Contact",9.2,2.9,2.3,0.35,22,GOLD,True)
textbox(s,"Ogunbunmi Philip Oluwadarasimi\nFederal University of Technology, Akure\n\n08169150242\nogunbunmiphilip@gmail.com",9.2,3.65,2.6,1.65,17,CREAM)
textbox(s,"SmartAqua",0.8,6.65,2.0,0.35,18,GOLD,True)
textbox(s,"Intelligent feeding for climate-resilient aquaculture",3.0,6.68,7.0,0.25,14,MUTED)

prs.core_properties.title = "SmartAqua - Africa Agri-future Innovation Challenge 2026"
prs.core_properties.author = "Ogunbunmi Philip Oluwadarasimi"
prs.core_properties.subject = "Pilot-stage intelligent aquaculture innovation"
prs.save(OUT)
print(OUT)
